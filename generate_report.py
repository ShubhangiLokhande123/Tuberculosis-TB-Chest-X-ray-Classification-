import os
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt

def generate_pdf():
    class PDF(FPDF):
        def header(self):
            self.set_font('Arial', 'B', 12)
            self.cell(0, 10, 'Tuberculosis Chest X-ray Classification Study Report', 0, 1, 'C')
            self.ln(5)

        def footer(self):
            self.set_y(-15)
            self.set_font('Arial', 'I', 8)
            self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')

        def chapter_title(self, title):
            self.set_font('Arial', 'B', 14)
            self.set_fill_color(200, 220, 255)
            self.cell(0, 10, title, 0, 1, 'L', 1)
            self.ln(4)

        def chapter_body(self, body):
            self.set_font('Arial', '', 11)
            self.multi_cell(0, 6, body)
            self.ln()

    pdf = PDF()
    pdf.add_page()

    # Title
    pdf.set_font('Arial', 'B', 16)
    pdf.cell(0, 15, 'Tuberculosis (TB) Chest X-ray Classification:', 0, 1, 'C')
    pdf.set_font('Arial', 'I', 14)
    pdf.cell(0, 10, 'Comparing Training from Scratch vs Transfer Learning', 0, 1, 'C')
    pdf.ln(10)

    # Introduction
    intro_text = (
        "Tuberculosis (TB) remains one of the top 10 causes of death worldwide according to the WHO. "
        "Early and accurate detection via chest X-ray is critical, especially in resource-limited "
        "settings where radiologist access is scarce. This study tackles the binary classification of "
        "chest X-rays, distinguishing Normal lungs from those showing signs of Tuberculosis (TB)."
    )
    pdf.chapter_title('1. Introduction')
    pdf.chapter_body(intro_text)

    # Objective
    objective_text = (
        "The primary objective of this project is to develop and evaluate deep learning models for automated "
        "TB detection in chest radiographs. Specifically, the study compares the performance of a custom "
        "Convolutional Neural Network (CNN) trained from scratch against four state-of-the-art architectures "
        "using transfer learning: ResNet50, DenseNet121, EfficientNet-B0, and Vision Transformer (ViT-B/16). "
        "A key clinical focus is placed on TB Recall (Sensitivity), as missing a TB case poses a severe "
        "public health risk."
    )
    pdf.chapter_title('2. Objective')
    pdf.chapter_body(objective_text)

    # Methodology
    methodology_text = (
        "The dataset used is the Tuberculosis (TB) Chest X-ray Dataset from Kaggle, consisting of ~4,200 "
        "images (~3,500 Normal and ~700 TB), creating a significant 5:1 class imbalance. "
        "Images were resized to 224x224 pixels. Data augmentation techniques such as random horizontal "
        "flips, rotations, color jitter, and affine translations were applied to the training set.\n\n"
        "Five models were developed:\n"
        "1. Custom CNN: A baseline 4-block ConvNet trained from scratch.\n"
        "2. ResNet50: A 50-layer residual network pre-trained on ImageNet.\n"
        "3. DenseNet121: A densely connected network pre-trained on ImageNet.\n"
        "4. EfficientNet-B0: A compound-scaled CNN pre-trained on ImageNet.\n"
        "5. ViT-B/16: A Vision Transformer pre-trained on ImageNet using patch-based self-attention.\n\n"
        "The data was split into 70% training, 15% validation, and 15% testing. For transfer learning models, "
        "differential learning rates were employed (backbone LR x 0.1, head LR x 1.0) to preserve low-level "
        "features while adapting the classification head. The Adam/AdamW optimizers were used with CrossEntropyLoss."
    )
    pdf.chapter_title('3. Methodology')
    pdf.chapter_body(methodology_text)

    # Results and Discussion
    results_text = (
        "The models were evaluated based on Test Accuracy, AUC-ROC, Macro F1, and TB Recall on the held-out "
        "test set of 630 images.\n\n"
        "Summary of Results:\n"
        "- Custom CNN (Scratch): Accuracy: ~88%, AUC: 0.928, TB Recall: ~73.3%\n"
        "- ResNet50: Accuracy: ~96.8%, AUC: 0.993, TB Recall: ~91.4%\n"
        "- DenseNet121: Accuracy: ~97.6%, AUC: 0.996, TB Recall: ~93.3%\n"
        "- EfficientNet-B0: Accuracy: ~97.9%, AUC: 0.997, TB Recall: ~95.2%\n"
        "- ViT-B/16: Accuracy: ~98.6%, AUC: 0.999, TB Recall: ~96.2%\n\n"
        "Discussion:\n"
        "Transfer learning models dramatically outperformed the baseline model trained from scratch, demonstrating "
        "the value of ImageNet pre-training even for specialized medical imaging tasks. The Vision Transformer "
        "(ViT-B/16) achieved the highest overall performance, demonstrating its ability to capture global long-range "
        "dependencies from the first layer, which is essential for interpreting whole-lung structural patterns."
    )
    pdf.chapter_title('4. Results and Discussion')
    pdf.chapter_body(results_text)

    # Conclusion
    conclusion_text = (
        "The study demonstrates that deep learning, particularly transfer learning and Vision Transformers, "
        "is highly effective for the binary classification of TB in chest X-rays. ViT-B/16 achieved the best "
        "clinical metrics, most notably the highest TB recall (sensitivity), minimizing the risk of false negatives. "
        "DenseNet121 and EfficientNet-B0 also offered excellent accuracy with fewer parameters, making them suitable "
        "for edge deployment. Future work may involve addressing class imbalance with SMOTE, exploring ensemble "
        "models, and enhancing interpretability using Grad-CAM or Attention Rollout."
    )
    pdf.chapter_title('5. Conclusion')
    pdf.chapter_body(conclusion_text)

    # References
    ref_text = (
        "1. Dosovitskiy, A., et al. (2020). An Image is Worth 16x16 Words: Transformers for Image Recognition at Scale. arXiv:2010.11929\n"
        "2. He, K., et al. (2015). Deep Residual Learning for Image Recognition. arXiv:1512.03385\n"
        "3. Huang, G., et al. (2016). Densely Connected Convolutional Networks (DenseNet). arXiv:1608.06993\n"
        "4. Tan, M. & Le, Q. (2019). EfficientNet: Rethinking Model Scaling for CNNs. arXiv:1905.11946\n"
        "5. Rahman, T., et al. Tuberculosis (TB) Chest X-ray Dataset. Kaggle."
    )
    pdf.chapter_title('6. References')
    pdf.chapter_body(ref_text)

    pdf.output('TB_Classification_Report.pdf')


def generate_pptx():
    prs = Presentation()

    # Title Slide (Slide 1)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Tuberculosis (TB) Chest X-ray Classification"
    subtitle.text = "Comparing Training from Scratch vs Transfer Learning\nCustom CNN, ResNet50, DenseNet121, EfficientNet-B0, ViT-B/16"

    # Introduction Slide (Slide 2)
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Introduction"
    tf = body.text_frame
    tf.text = "Tuberculosis is a top 10 cause of death worldwide."
    p = tf.add_paragraph()
    p.text = "Early, accurate detection via chest X-ray is critical, especially in resource-limited settings."
    p = tf.add_paragraph()
    p.text = "This project aims to automate binary classification: Normal vs Tuberculosis."

    # Objective Slide (Slide 3)
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Objective"
    tf = body.text_frame
    tf.text = "Compare training deep learning models from scratch against transfer learning approaches."
    p = tf.add_paragraph()
    p.text = "Evaluate five distinct architectures: Custom CNN, ResNet50, DenseNet121, EfficientNet-B0, ViT-B/16."
    p = tf.add_paragraph()
    p.text = "Primary clinical focus: TB Recall (Sensitivity) to minimize false negatives."

    # Dataset & Preprocessing Slide (Slide 4)
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Dataset & Preprocessing"
    tf = body.text_frame
    tf.text = "Kaggle Tuberculosis (TB) Chest X-ray Dataset."
    p = tf.add_paragraph()
    p.text = "~4,200 images total (~3,500 Normal, ~700 TB) - class imbalance addressed."
    p = tf.add_paragraph()
    p.text = "Images resized to 224x224."
    p = tf.add_paragraph()
    p.text = "Data Augmentation: Random horizontal flips, rotation, color jitter, affine translations."

    # Methodology Slide (Slide 5)
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Methodology"
    tf = body.text_frame
    tf.text = "Data Splits: 70% Train, 15% Validation, 15% Test."
    p = tf.add_paragraph()
    p.text = "Baseline: 4-block Custom CNN trained from scratch."
    p = tf.add_paragraph()
    p.text = "Transfer Learning Models: ResNet50, DenseNet121, EfficientNet-B0, ViT-B/16 pre-trained on ImageNet."
    p = tf.add_paragraph()
    p.text = "Differential Learning Rates used for transfer learning: low LR for backbone, higher LR for head."

    # Model Architectures Slide (Slide 6)
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Model Architectures"
    tf = body.text_frame
    tf.text = "Custom CNN: Local receptive field, builds context hierarchically."
    p = tf.add_paragraph()
    p.text = "DenseNet121: Dense feature reuse suited to medical imaging."
    p = tf.add_paragraph()
    p.text = "EfficientNet-B0: Efficient compound-scaled CNN."
    p = tf.add_paragraph()
    p.text = "ViT-B/16 (Vision Transformer): Global receptive field from layer 1 via patch-based self-attention."

    # Results Slide (Slide 7)
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Results"
    tf = body.text_frame
    tf.text = "Custom CNN: ~88% Acc, 0.928 AUC, ~73% TB Recall"
    p = tf.add_paragraph()
    p.text = "ResNet50: ~96.8% Acc, 0.993 AUC, ~91.4% TB Recall"
    p = tf.add_paragraph()
    p.text = "DenseNet121: ~97.6% Acc, 0.996 AUC, ~93.3% TB Recall"
    p = tf.add_paragraph()
    p.text = "EfficientNet-B0: ~97.9% Acc, 0.997 AUC, ~95.2% TB Recall"
    p = tf.add_paragraph()
    p.text = "ViT-B/16: ~98.6% Acc, 0.999 AUC, ~96.2% TB Recall"

    # Discussion Slide (Slide 8)
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Discussion"
    tf = body.text_frame
    tf.text = "Transfer learning dramatically outperforms training from scratch."
    p = tf.add_paragraph()
    p.text = "ImageNet pre-training covers the requirement for massive data."
    p = tf.add_paragraph()
    p.text = "ViT-B/16 achieves the highest accuracy, AUC-ROC, and crucial TB Recall, making it ideal for clinical screening."

    # Conclusion Slide (Slide 9)
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Conclusion & Recommendations"
    tf = body.text_frame
    tf.text = "Vision Transformers excel in medical image classification tasks due to their global long-range dependencies."
    p = tf.add_paragraph()
    p.text = "Deployment Recommendations: ViT-B/16 for maximum accuracy; EfficientNet-B0 for constrained compute."
    p = tf.add_paragraph()
    p.text = "Future Work: Incorporate Attention Rollout/Grad-CAM for interpretability and SMOTE for class imbalance."

    # References Slide (Slide 10)
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "References"
    tf = body.text_frame
    tf.text = "Dosovitskiy, A., et al. (2020). An Image is Worth 16x16 Words."
    p = tf.add_paragraph()
    p.text = "He, K., et al. (2015). Deep Residual Learning for Image Recognition."
    p = tf.add_paragraph()
    p.text = "Huang, G., et al. (2016). Densely Connected Convolutional Networks."
    p = tf.add_paragraph()
    p.text = "Tan, M. & Le, Q. (2019). EfficientNet: Rethinking Model Scaling."

    prs.save('TB_Classification_Presentation.pptx')


if __name__ == '__main__':
    generate_pdf()
    generate_pptx()
    print("Report and Presentation generated successfully!")
